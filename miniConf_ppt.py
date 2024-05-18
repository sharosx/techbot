import sys
import warnings
import os
import constants
from langchain_community.document_loaders import PyPDFLoader
from langchain.indexes import VectorstoreIndexCreator
import ssl
import socket
import streamlit as st
from pptx import Presentation

# Suppress all warnings
warnings.filterwarnings("ignore")

# Set environment variable
os.environ["OPENAI_API_KEY"] = constants.API_KEY

# Get query from command-line argument
if len(sys.argv) > 1:
    query = sys.argv[1]
else:
    query = "adasis"  # Default query if not provided

# Load PDF and PowerPoint files
pdf_file_path = "/Users/sharosmassaebi/PycharmProjects/pythonProject/Adasis_v2.pdf"
pptx_file_path = "AI Data Academy.pptx"

if os.path.exists(pdf_file_path) and os.path.exists(pptx_file_path):
    # Load PDF file
    pdf_loader = PyPDFLoader(pdf_file_path)
    pdf_content = pdf_loader.load()

    # Extract text from PowerPoint file
    ppt_text = []
    prs = Presentation(pptx_file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text.append(shape.text.strip())

    # Combine PDF and PowerPoint text
    combined_text = "\n".join(ppt_text)

    # Create vector store index
    index_creator = VectorstoreIndexCreator()

    # Add PDF content
    index_creator.add_document("pdf_document", pdf_content)

    # Add PowerPoint text
    index_creator.add_document("pptx_document", combined_text)

    # Create index
    index = index_creator.create_index()

    # Query index
    try:
        result = index.query(query)
        if result:
            st.write("Result:", result)
        else:
            st.write("No result found for the query.")
    except Exception as e:
        st.error(f"An error occurred: {e}")
    finally:
        # Explicitly close the SSL context after use to avoid ResourceWarning
        context = ssl.create_default_context()
        context.check_hostname = False
        context.verify_mode = ssl.CERT_NONE
        ssl_socket = context.wrap_socket(socket.socket(socket.AF_INET), suppress_ragged_eofs=True, server_hostname="")
        ssl_socket.close()

        # Reset warnings after executing potentially problematic code
        warnings.resetwarnings()
else:
    st.error("One or both of the specified files do not exist.")
